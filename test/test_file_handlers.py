import io
import os
from unittest.mock import patch, MagicMock

import pytest

from backend.file_handlers import (
    extract_text_from_docx,
    extract_text_from_xlsx,
    extract_text_from_pdf,
    extract_text_from_pptx,
    anonymize_docx,
    anonymize_xlsx,
    anonymize_pdf,
    anonymize_pptx,
    restore_docx,
    restore_xlsx,
    restore_pdf,
    restore_pptx,
)


class TestExtractTextDocx:
    def test_extract_docx_text(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Hello World")
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        text, loaded_doc = extract_text_from_docx(buf.read())
        assert "Hello World" in text


class TestAnonymizeDocx:
    def test_anonymize_replaces_entities(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Meeting with John Doe")
        entity_map = {"PERSON_1": "John Doe"}
        counts, number_map = anonymize_docx(doc, entity_map, 1.0)
        assert counts["entities"] > 0
        assert "John Doe" not in doc.paragraphs[0].text
        assert "PERSON_1" in doc.paragraphs[0].text

    def test_anonymize_number_multiplier(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Amount 1000 dollars")
        entity_map = {}
        counts, number_map = anonymize_docx(doc, entity_map, 2.0)
        assert counts["numbers"] > 0

    def test_anonymize_no_multiplier(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Amount 1000 dollars")
        entity_map = {}
        counts, number_map = anonymize_docx(doc, entity_map, 1.0)
        assert counts["numbers"] == 0

    def test_restore_docx(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Meeting with PERSON_1")
        key_data = {"entities": {"PERSON_1": "John Doe"}, "number_mappings": {}, "multiplier": 1.0}
        counts = restore_docx(doc, key_data)
        assert counts["entities"] > 0
        assert "PERSON_1" not in doc.paragraphs[0].text
        assert "John Doe" in doc.paragraphs[0].text


class TestAnonymizeXLSX:
    def test_anonymize_string_cell(self):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "John Doe"
        entity_map = {"PERSON_1": "John Doe"}
        counts, number_map = anonymize_xlsx(wb, entity_map, 1.0)
        assert ws["A1"].value == "PERSON_1"

    def test_anonymize_number_cell(self):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = 1000
        entity_map = {}
        counts, number_map = anonymize_xlsx(wb, entity_map, 2.0)
        assert ws["A1"].value == 2000

    def test_restore_xlsx(self):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "PERSON_1"
        key_data = {"entities": {"PERSON_1": "John Doe"}, "number_mappings": {}, "multiplier": 1.0}
        counts = restore_xlsx(wb, key_data)
        assert ws["A1"].value == "John Doe"


class TestAnonymizePPTX:
    def test_anonymize_pptx_text(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "John Doe"
        entity_map = {"PERSON_1": "John Doe"}
        counts, number_map = anonymize_pptx(prs, entity_map, 1.0)
        assert tf.text == "PERSON_1"

    def test_restore_pptx(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "PERSON_1"
        key_data = {"entities": {"PERSON_1": "John Doe"}, "number_mappings": {}, "multiplier": 1.0}
        counts = restore_pptx(prs, key_data)
        assert tf.text == "John Doe"


class TestAnonymizePDF:
    def test_anonymize_text_pdf(self):
        import fitz

        pdf = fitz.open()
        page = pdf.new_page()
        tw = fitz.TextWriter(page.rect)
        tw.append(fitz.Point(50, 50), "John Doe")
        tw.write_text(page)
        buf = io.BytesIO()
        pdf.save(buf)
        pdf.close()
        buf.seek(0)

        entity_map = {"PERSON_1": "John Doe"}
        counts, pdf_bytes, number_map = anonymize_pdf(buf.read(), entity_map, 1.0)
        assert counts["entities"] > 0

    def test_restore_pdf(self):
        import fitz

        pdf = fitz.open()
        page = pdf.new_page()
        tw = fitz.TextWriter(page.rect)
        tw.append(fitz.Point(50, 50), "PERSON_1")
        tw.write_text(page)
        buf = io.BytesIO()
        pdf.save(buf)
        pdf.close()
        buf.seek(0)

        key_data = {"entities": {"PERSON_1": "John Doe"}, "number_mappings": {}, "multiplier": 1.0}
        counts, pdf_bytes = restore_pdf(buf.read(), key_data)
        assert counts["entities"] > 0


class TestExtractTextPDF:
    def test_extract_text_from_text_pdf(self):
        import fitz

        pdf = fitz.open()
        page = pdf.new_page()
        tw = fitz.TextWriter(page.rect)
        tw.append(fitz.Point(50, 50), "Hello World")
        tw.write_text(page)
        buf = io.BytesIO()
        pdf.save(buf)
        pdf.close()
        buf.seek(0)

        text, pages = extract_text_from_pdf(buf.read())
        assert "Hello World" in text
