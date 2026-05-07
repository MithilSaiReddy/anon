import io
import re
from typing import Tuple, Dict, Any, List

from docx import Document
from docx.text.paragraph import Paragraph
import openpyxl
import pdfplumber
import fitz
from pptx import Presentation
from pptx.util import Emu


def extract_text_from_docx(file_bytes: bytes) -> Tuple[str, Document]:
    doc = Document(io.BytesIO(file_bytes))
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text.append(cell.text)
    return "\n".join(full_text), doc


def extract_text_from_xlsx(file_bytes: bytes) -> Tuple[str, openpyxl.Workbook]:
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    full_text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    full_text.append(str(cell))
    return "\n".join(full_text), wb


def extract_text_from_pdf(file_bytes: bytes) -> Tuple[str, List[str]]:
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        pages = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n".join(pages), pages


def _replace_entities_in_text(text: str, entity_map: Dict[str, str], multiplier: float, number_map: Dict[str, float]) -> Tuple[str, int]:
    count = 0
    for placeholder, original in entity_map.items():
        if original in text:
            text = text.replace(original, placeholder)
            count += 1

    if multiplier and multiplier != 1.0:
        def multiply_number(match):
            num_str = match.group()
            try:
                num = float(num_str)
                new_num = round(num * multiplier, 2)
                number_map[num_str] = new_num
                if new_num == int(new_num):
                    return str(int(new_num))
                return str(new_num)
            except ValueError:
                return num_str

        text = re.sub(r'\b\d+\.?\d+\b', multiply_number, text)

    return text, count


def anonymize_docx(doc: Document, entity_map: Dict[str, str], multiplier: float = 1.0) -> Dict[str, Any]:
    number_map = {}
    total_replacements = 0

    def process_run(run):
        nonlocal total_replacements
        if not run.text:
            return
        text = run.text
        for placeholder, original in entity_map.items():
            if original in text:
                text = text.replace(original, placeholder)
                total_replacements += 1

        if multiplier and multiplier != 1.0:
            def multiply_number(match):
                num_str = match.group()
                try:
                    num = float(num_str)
                    new_num = round(num * multiplier, 2)
                    number_map[num_str] = new_num
                    if new_num == int(new_num):
                        return str(int(new_num))
                    return str(new_num)
                except ValueError:
                    return num_str

            new_text, _count = re.subn(r'\b\d+\.?\d+\b', multiply_number, text)
            if _count > 0:
                text = new_text
                total_replacements += _count

        run.text = text

    for para in doc.paragraphs:
        for run in para.runs:
            process_run(run)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        process_run(run)

    return {"entities": total_replacements, "numbers": len(number_map)}, number_map


def anonymize_xlsx(wb: openpyxl.Workbook, entity_map: Dict[str, str], multiplier: float = 1.0) -> Dict[str, Any]:
    from openpyxl.cell.cell import MergedCell

    number_map = {}
    total_replacements = 0

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                try:
                    if isinstance(cell, MergedCell):
                        continue
                except Exception:
                    pass

                if cell.value is None:
                    continue

                if isinstance(cell.value, str):
                    text = cell.value
                    for placeholder, original in entity_map.items():
                        if original in text:
                            text = text.replace(original, placeholder)
                            total_replacements += 1

                    if multiplier and multiplier != 1.0:
                        def multiply_number(match):
                            num_str = match.group()
                            try:
                                num = float(num_str)
                                new_num = round(num * multiplier, 2)
                                number_map[num_str] = new_num
                                if new_num == int(new_num):
                                    return str(int(new_num))
                                return str(new_num)
                            except ValueError:
                                return num_str

                        new_text, _count = re.subn(r'\b\d+\.?\d+\b', multiply_number, text)
                        if _count > 0:
                            text = new_text
                            total_replacements += _count

                    if text != cell.value:
                        try:
                            cell.value = text
                        except Exception:
                            pass

                elif isinstance(cell.value, (int, float)):
                    if multiplier and multiplier != 1.0:
                        original_val = cell.value
                        new_val = round(original_val * multiplier, 2)
                        number_map[str(original_val)] = new_val
                        cell.value = new_val

    return {"entities": total_replacements, "numbers": len(number_map)}, number_map


def extract_text_from_pptx(file_bytes: bytes) -> Tuple[str, Presentation]:
    prs = Presentation(io.BytesIO(file_bytes))
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            full_text.append(para.text)
    return "\n".join(full_text), prs


def anonymize_pptx(prs: Presentation, entity_map: Dict[str, str], multiplier: float = 1.0) -> Dict[str, Any]:
    number_map = {}
    total_replacements = 0

    def process_text(text: str) -> str:
        nonlocal total_replacements
        for placeholder, original in entity_map.items():
            if original in text:
                text = text.replace(original, placeholder)
                total_replacements += 1

        if multiplier and multiplier != 1.0:
            def multiply_number(match):
                num_str = match.group()
                try:
                    num = float(num_str)
                    new_num = round(num * multiplier, 2)
                    number_map[num_str] = new_num
                    if new_num == int(new_num):
                        return str(int(new_num))
                    return str(new_num)
                except ValueError:
                    return num_str

            text = re.sub(r'\b\d+\.?\d+\b', multiply_number, text)
        return text

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = process_text(run.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = process_text(run.text)

    return {"entities": total_replacements, "numbers": len(number_map)}, number_map


def restore_pptx(prs: Presentation, key_data: Dict[str, Any]) -> Dict[str, int]:
    entity_map = key_data.get("entities", {})
    number_map = key_data.get("number_mappings", {})
    multiplier = key_data.get("multiplier", 1.0)

    restore_count = {"entities": 0, "numbers": 0}

    def process_text(text: str) -> str:
        for placeholder, original in entity_map.items():
            if placeholder in text:
                text = text.replace(placeholder, original)
                restore_count["entities"] += 1

        if number_map:
            for orig_str, new_val in number_map.items():
                if str(new_val) in text:
                    text = text.replace(str(new_val), orig_str)
                    restore_count["numbers"] += 1
        elif multiplier != 1.0:
            def divide_number(match):
                try:
                    num = float(match.group())
                    result = round(num / multiplier, 2)
                    if result == int(result):
                        return str(int(result))
                    return str(result)
                except ValueError:
                    return match.group()

            text = re.sub(r'\b\d+\.?\d+\b', divide_number, text)
        return text

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = process_text(run.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = process_text(run.text)

    return restore_count


def anonymize_pdf(file_bytes: bytes, entity_map: Dict[str, str], multiplier: float = 1.0) -> Tuple[Dict[str, Any], bytes]:
    number_map = {}
    total_replacements = 0

    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    all_replacements = []

    for placeholder, original in entity_map.items():
        if not original or len(original) < 2:
            continue
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            instances = page.search_for(original)
            for inst in instances:
                all_replacements.append((page_num, inst, original, placeholder))

    if multiplier and multiplier != 1.0:
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            words = page.get_text("words")
            for word_info in words:
                x0, y0, x1, y1, word_text, block_no, line_no, word_no = word_info
                if re.match(r'^\d+\.?\d+$', word_text):
                    try:
                        num = float(word_text)
                        new_num = round(num * multiplier, 2)
                        number_map[word_text] = new_num
                        if new_num == int(new_num):
                            new_text = str(int(new_num))
                        else:
                            new_text = str(new_num)
                        rect = fitz.Rect(x0, y0, x1, y1)
                        all_replacements.append((page_num, rect, word_text, new_text))
                    except ValueError:
                        pass

    processed_positions = set()

    for page_num, rect, old_text, new_text in all_replacements:
        pos_key = (page_num, round(rect.x0, 1), round(rect.y0, 1))
        if pos_key in processed_positions:
            continue
        processed_positions.add(pos_key)

        page = pdf[page_num]
        page.add_redact_annot(rect, text=new_text, fill=(1, 1, 1), fontsize=None)
        total_replacements += 1

    for page_num in range(len(pdf)):
        pdf[page_num].apply_redactions()

    output_buffer = io.BytesIO()
    pdf.save(output_buffer)
    output_buffer.seek(0)
    pdf_bytes = output_buffer.getvalue()
    pdf.close()

    return {"entities": total_replacements, "numbers": len(number_map)}, pdf_bytes, number_map


def restore_pdf(file_bytes: bytes, key_data: Dict[str, Any]) -> Tuple[Dict[str, int], bytes]:
    entity_map = key_data.get("entities", {})
    number_map = key_data.get("number_mappings", {})
    multiplier = key_data.get("multiplier", 1.0)

    restore_count = {"entities": 0, "numbers": 0}

    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    all_replacements = []

    for placeholder, original in entity_map.items():
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            instances = page.search_for(placeholder)
            for inst in instances:
                all_replacements.append((page_num, inst, placeholder, original))

    if number_map:
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            words = page.get_text("words")
            for word_info in words:
                x0, y0, x1, y1, word_text, _, _, _ = word_info
                for orig_str, new_val in number_map.items():
                    if word_text.replace(",", "") == str(new_val).replace(",", ""):
                        rect = fitz.Rect(x0, y0, x1, y1)
                        all_replacements.append((page_num, rect, word_text, orig_str))
    elif multiplier != 1.0:
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            words = page.get_text("words")
            for word_info in words:
                x0, y0, x1, y1, word_text, _, _, _ = word_info
                if re.match(r'^\d+\.?\d+$', word_text):
                    try:
                        num = float(word_text)
                        result = round(num / multiplier, 2)
                        new_text = str(int(result)) if result == int(result) else str(result)
                        rect = fitz.Rect(x0, y0, x1, y1)
                        all_replacements.append((page_num, rect, word_text, new_text))
                    except ValueError:
                        pass

    processed_positions = set()

    for page_num, rect, old_text, new_text in all_replacements:
        pos_key = (page_num, round(rect.x0, 1), round(rect.y0, 1))
        if pos_key in processed_positions:
            continue
        processed_positions.add(pos_key)

        page = pdf[page_num]
        page.add_redact_annot(rect, text=new_text, fill=(1, 1, 1), fontsize=None)
        if old_text in entity_map:
            restore_count["entities"] += 1
        else:
            restore_count["numbers"] += 1

    for page_num in range(len(pdf)):
        pdf[page_num].apply_redactions()

    output_buffer = io.BytesIO()
    pdf.save(output_buffer)
    output_buffer.seek(0)
    pdf_bytes = output_buffer.getvalue()
    pdf.close()

    return restore_count, pdf_bytes


def create_pdf_from_text(pages: List[str]) -> bytes:
    pdf = fitz.open()
    for page_text in pages:
        page = pdf.new_page()
        rect = page.rect
        tw = fitz.TextWriter(rect)
        y = 50
        for line in page_text.split('\n'):
            if y > rect.height - 30:
                break
            tw.append((50, y), line)
            y += 14
        tw.write_text(page)

    output_buffer = io.BytesIO()
    pdf.save(output_buffer)
    output_buffer.seek(0)
    return output_buffer.getvalue()


def restore_docx(doc: Document, key_data: Dict[str, Any]) -> Dict[str, int]:
    entity_map = key_data.get("entities", {})
    number_map = key_data.get("number_mappings", {})
    multiplier = key_data.get("multiplier", 1.0)

    reverse_entity = {v: k for k, v in entity_map.items()}
    restore_count = {"entities": 0, "numbers": 0}

    def process_run(run):
        if not run.text:
            return
        text = run.text

        for placeholder, original in entity_map.items():
            if placeholder in text:
                text = text.replace(placeholder, original)
                restore_count["entities"] += 1

        if multiplier != 1.0:
            for orig_str, new_val in number_map.items():
                if str(new_val) in text:
                    text = text.replace(str(new_val), orig_str)
                    restore_count["numbers"] += 1
            if not number_map:
                def divide_number(match):
                    try:
                        num = float(match.group())
                        result = round(num / multiplier, 2)
                        if result == int(result):
                            return str(int(result))
                        return str(result)
                    except ValueError:
                        return match.group()

                text = re.sub(r'\b\d+\.?\d+\b', divide_number, text)

        run.text = text

    for para in doc.paragraphs:
        for run in para.runs:
            process_run(run)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        process_run(run)

    return restore_count


def restore_xlsx(wb: openpyxl.Workbook, key_data: Dict[str, Any]) -> Dict[str, int]:
    from openpyxl.cell.cell import MergedCell

    entity_map = key_data.get("entities", {})
    number_map = key_data.get("number_mappings", {})
    multiplier = key_data.get("multiplier", 1.0)

    restore_count = {"entities": 0, "numbers": 0}

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                try:
                    if isinstance(cell, MergedCell):
                        continue
                except Exception:
                    pass

                if cell.value is None:
                    continue

                if isinstance(cell.value, str):
                    text = cell.value
                    for placeholder, original in entity_map.items():
                        if placeholder in text:
                            text = text.replace(placeholder, original)
                            restore_count["entities"] += 1

                    if number_map:
                        for orig_str, new_val in number_map.items():
                            if str(new_val) in text:
                                text = text.replace(str(new_val), orig_str)
                                restore_count["numbers"] += 1
                    elif multiplier != 1.0:
                        def divide_number(match):
                            try:
                                num = float(match.group())
                                result = round(num / multiplier, 2)
                                if result == int(result):
                                    return str(int(result))
                                return str(result)
                            except ValueError:
                                return match.group()

                        text = re.sub(r'\b\d+\.?\d+\b', divide_number, text)

                    if text != cell.value:
                        try:
                            cell.value = text
                        except Exception:
                            pass

                elif isinstance(cell.value, (int, float)):
                    if number_map:
                        for orig_str, new_val in number_map.items():
                            try:
                                if abs(cell.value - float(new_val)) < 0.001:
                                    cell.value = float(orig_str)
                                    restore_count["numbers"] += 1
                                    break
                            except (ValueError, TypeError):
                                pass
                    elif multiplier != 1.0:
                        original_val = cell.value
                        cell.value = round(original_val / multiplier, 2)
                        restore_count["numbers"] += 1

    return restore_count


def restore_pdf_text(pages: List[str], key_data: Dict[str, Any]) -> Tuple[Dict[str, int], List[str]]:
    entity_map = key_data.get("entities", {})
    number_map = key_data.get("number_mappings", {})
    multiplier = key_data.get("multiplier", 1.0)

    restore_count = {"entities": 0, "numbers": 0}
    new_pages = []

    for page_text in pages:
        text = page_text

        for placeholder, original in entity_map.items():
            if placeholder in text:
                text = text.replace(placeholder, original)
                restore_count["entities"] += 1

        if number_map:
            for orig_str, new_val in number_map.items():
                if str(new_val) in text:
                    text = text.replace(str(new_val), orig_str)
                    restore_count["numbers"] += 1
        elif multiplier != 1.0:
            def divide_number(match):
                try:
                    num = float(match.group())
                    result = round(num / multiplier, 2)
                    if result == int(result):
                        return str(int(result))
                    return str(result)
                except ValueError:
                    return match.group()

            text = re.sub(r'\b\d+\.?\d+\b', divide_number, text)

        new_pages.append(text)

    return restore_count, new_pages
